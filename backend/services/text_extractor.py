"""
文本提取服务
支持从 PDF、DOCX、PPTX、TXT 文件中提取文本
支持 OCR 识别扫描件
"""
import os
import io
from typing import Optional
import pypdfium2 as pdfium
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

try:
    from google.cloud import vision
    GOOGLE_VISION_AVAILABLE = True
except ImportError:
    GOOGLE_VISION_AVAILABLE = False


class TextExtractor:
    """文本提取器类"""
    
    def __init__(self):
        """初始化文本提取器"""
        self.use_google_vision = (
            GOOGLE_VISION_AVAILABLE and 
            os.getenv('GOOGLE_APPLICATION_CREDENTIALS')
        )
    
    def extract_text(self, file_path: str, file_extension: str) -> str:
        """
        根据文件类型提取文本
        
        Args:
            file_path: 文件路径
            file_extension: 文件扩展名 (包含点，如 '.pdf')
        
        Returns:
            提取的文本内容
        
        Raises:
            ValueError: 不支持的文件格式
        """
        extension = file_extension.lower()
        
        if extension == '.pdf':
            return self._extract_from_pdf(file_path)
        elif extension == '.docx':
            return self._extract_from_docx(file_path)
        elif extension == '.pptx':
            return self._extract_from_pptx(file_path)
        elif extension == '.txt':
            return self._extract_from_txt(file_path)
        else:
            raise ValueError(f"不支持的文件格式: {extension}")
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """从 PDF 文件提取文本"""
        text_content = []
        
        try:
            pdf = pdfium.PdfDocument(file_path)
            
            for page_num in range(len(pdf)):
                page = pdf[page_num]
                textpage = page.get_textpage()
                text = textpage.get_text_range()
                
                # 如果页面没有文本，尝试 OCR
                if not text.strip():
                    # 渲染页面为图像
                    pil_image = page.render(
                        scale=2,
                        rotation=0,
                    ).to_pil()
                    text = self._ocr_image(pil_image)
                
                text_content.append(text)
            
            pdf.close()
            return '\n\n'.join(text_content)
        
        except Exception as e:
            raise Exception(f"PDF 处理失败: {str(e)}")
    
    def _extract_from_docx(self, file_path: str) -> str:
        """从 DOCX 文件提取文本"""
        try:
            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return '\n\n'.join(paragraphs)
        except Exception as e:
            raise Exception(f"DOCX 处理失败: {str(e)}")
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """从 PPTX 文件提取文本"""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_content.append('\n'.join(slide_text))
            
            return '\n\n'.join(text_content)
        except Exception as e:
            raise Exception(f"PPTX 处理失败: {str(e)}")
    
    def _extract_from_txt(self, file_path: str) -> str:
        """从 TXT 文件提取文本"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    return f.read()
            except Exception as e:
                raise Exception(f"TXT 处理失败: {str(e)}")
        except Exception as e:
            raise Exception(f"TXT 处理失败: {str(e)}")
    
    def _ocr_image(self, image: Image.Image) -> str:
        """
        对图像进行 OCR 识别
        
        Args:
            image: PIL Image 对象
        
        Returns:
            识别的文本
        """
        if self.use_google_vision:
            return self._google_vision_ocr(image)
        else:
            return self._tesseract_ocr(image)
    
    def _tesseract_ocr(self, image: Image.Image) -> str:
        """使用 Tesseract 进行 OCR"""
        try:
            # 支持中英文
            text = pytesseract.image_to_string(
                image, 
                lang='chi_sim+eng',
                config='--psm 3'
            )
            return text
        except Exception as e:
            print(f"Tesseract OCR 失败: {str(e)}")
            return ""
    
    def _google_vision_ocr(self, image: Image.Image) -> str:
        """使用 Google Cloud Vision 进行 OCR"""
        try:
            client = vision.ImageAnnotatorClient()
            
            # 将 PIL Image 转换为字节
            img_byte_arr = io.BytesIO()
            image.save(img_byte_arr, format='PNG')
            img_byte_arr = img_byte_arr.getvalue()
            
            vision_image = vision.Image(content=img_byte_arr)
            response = client.text_detection(image=vision_image)
            
            if response.text_annotations:
                return response.text_annotations[0].description
            return ""
        
        except Exception as e:
            print(f"Google Vision OCR 失败: {str(e)}")
            # 回退到 Tesseract
            return self._tesseract_ocr(image)
