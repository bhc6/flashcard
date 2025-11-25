import os
import PyPDF2
from openai import OpenAI, APIConnectionError
from dotenv import load_dotenv
import json
import hashlib
import unicodedata

# Language setting: 'en' for English, 'zh' for Chinese
LANGUAGE = 'zh'  # 可选: 'en' 或 'zh'

# Maximum number of flashcards to generate (set to 0 for unlimited)
MAX_FLASHCARDS = 0  # 0表示无限制，其他正整数表示最大生成数量

try:
    import pytesseract
    from PIL import Image
    from pdf2image import convert_from_path
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

ROOT_DIRECTORY = os.path.dirname(os.path.realpath(__file__))
load_dotenv()
ARK_BASE_URL = "https://ark.cn-beijing.volces.com/api/v3"

def get_openai_client():
    api_key = os.environ.get("ARK_API_KEY")
    if not api_key:
        raise SystemExit("Set ARK_API_KEY before running this script.")
    return OpenAI(api_key=api_key, base_url=ARK_BASE_URL)

client = get_openai_client()

# Read PDF with OCR support for scanned documents
def read_pdf(file_path):
    text = ""
    
    # First, try to extract text directly
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + " "
    
    # If no text extracted and OCR is available, use OCR on images
    if not text.strip() and OCR_AVAILABLE:
        print("No text found. Attempting OCR on scanned document...")
        try:
            images = convert_from_path(file_path)
            for img in images:
                ocr_text = pytesseract.image_to_string(img, lang='eng+chi_sim')
                text += ocr_text + " "
            print(f"OCR extraction completed: {len(text)} characters extracted")
        except Exception as e:
            print(f"OCR processing failed: {e}")
    elif not text.strip():
        print("Warning: No text extracted and OCR not available. Install: pip install pytesseract pdf2image pillow")
    
    return text

# Read PowerPoint file
def read_pptx(file_path):
    text = ""
    try:
        presentation = Presentation(file_path)
        for slide_num, slide in enumerate(presentation.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
        print(f"PPTX extraction completed: {len(text)} characters extracted from {len(presentation.slides)} slides")
    except Exception as e:
        print(f"PPTX processing failed: {e}")
    return text

# Read PowerPoint 97-2003 file (.ppt)
def read_ppt(file_path):
    text = ""
    try:
        import subprocess
        import tempfile
        
        # Try conversion method 1: unoconv
        try:
            subprocess.run(['unoconv', '-f', 'pptx', '-o', file_path.replace('.ppt', '_converted.pptx'), file_path], 
                          check=True, capture_output=True, timeout=30)
            converted_file = file_path.replace('.ppt', '_converted.pptx')
            if os.path.exists(converted_file):
                text = read_pptx(converted_file)
                os.remove(converted_file)
                print(f"PPT conversion successful using unoconv")
                return text
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
        
        # Try conversion method 2: libreoffice soffice
        try:
            temp_dir = tempfile.gettempdir()
            subprocess.run(['soffice', '--headless', '--convert-to', 'pptx', '--outdir', temp_dir, file_path],
                          check=True, capture_output=True, timeout=30)
            converted_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(file_path))[0] + '.pptx')
            if os.path.exists(converted_file):
                text = read_pptx(converted_file)
                os.remove(converted_file)
                print(f"PPT conversion successful using soffice")
                return text
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
        
        # Fallback: Extract using zipfile (PPT files are ZIP archives with OLE streams)
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                # Try to find slide files
                slide_files = [f for f in zip_ref.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                
                if slide_files:
                    for slide_file in sorted(slide_files):
                        try:
                            with zip_ref.open(slide_file) as xml_file:
                                root = ET.fromstring(xml_file.read())
                                # Extract all text elements from the slide
                                for elem in root.iter():
                                    if elem.text:
                                        text += elem.text + " "
                        except Exception:
                            pass
                    
                    if text.strip():
                        print(f"PPT extraction completed using zipfile: {len(text)} characters extracted")
                        return text
        except (zipfile.BadZipFile, Exception):
            pass
        
        print("Warning: Could not parse PPT file.")
        print("Recommended solutions:")
        print("  1. Convert PPT to PPTX manually and use the PPTX file")
        print("  2. Install unoconv: pip install unoconv (Linux/Mac)")
        print("  3. Install LibreOffice: sudo apt-get install libreoffice (Linux)")
    
    except Exception as e:
        print(f"PPT processing failed: {e}")
    
    return text

# dividing text into smaller chunks:
def divide_text(text, section_size):
    sections = []
    start = 0
    end = section_size
    while start < len(text):
        section = text[start:end]
        sections.append(section)
        start = end
        end += section_size
    return sections

# Language prompts
PROMPTS = {
    'en': {
        'system': "You are a helpful assistant.",
        'user': "Create anki flashcards with the provided text using a format: question;answer newline question;answer etc. Keep question and the corresponding answer on the same line. Do not add any introductory text. Text: {text}"
    },
    'zh': {
        'system': "你是一个有帮助的助手。",
        'user': "根据提供的文本创建Anki闪卡，使用以下格式：问题;答案 换行 问题;答案 等。确保问题和对应的答案在同一行。不要添加任何介绍文本。文本：{text}"
    }
}

# Convert filename to ASCII-compatible name
def sanitize_filename(filename):
    """Convert filename to ASCII-safe format for JSON output"""
    base_name = os.path.splitext(filename)[0]
    
    # Try to normalize and remove accents
    try:
        # Check if filename contains non-ASCII characters
        base_name.encode('ascii')
        return base_name
    except UnicodeEncodeError:
        # Generate hash-based name for non-ASCII filenames
        hash_suffix = hashlib.md5(base_name.encode('utf-8')).hexdigest()[:8]
        return f"flashcard_{hash_suffix}"

# Create Anki cards and save as JSON
def create_anki_cards(pdf_text, pdf_filename="document"):
    SECTION_SIZE = 1000
    divided_sections = divide_text(pdf_text, SECTION_SIZE)
    
    flashcards_list = []
    
    # 限制API调用次数，当前只处理第一部分文本
    for i, text in enumerate(divided_sections):
        if i > 0:
            break

        # Get prompts based on language setting
        lang = LANGUAGE if LANGUAGE in PROMPTS else 'en'
        prompts = PROMPTS[lang]
        
        messages = [
            {"role": "system", "content": prompts['system']},
            {"role": "user", "content": prompts['user'].format(text=text)}
        ]

        try:
            response = client.chat.completions.create(
                model="doubao-seed-1-6-251015",
                messages=messages,
                temperature=0.3,
                max_tokens=2048,
            )
            response_from_api = response.choices[0].message.content.strip()
            
            # 解析API返回的字符串
            lines = response_from_api.split('\n')
            for line in lines:
                # Check if max flashcards limit reached
                if MAX_FLASHCARDS > 0 and len(flashcards_list) >= MAX_FLASHCARDS:
                    break
                
                if ';' in line:
                    parts = line.split(';', 1)
                    question = parts[0].strip()
                    answer = parts[1].strip()
                    if question and answer:
                        flashcards_list.append({"question": question, "answer": answer})

        except APIConnectionError as exc:
            print(f"Connection error while calling the API: {exc}")
            continue
        except IndexError:
            print(f"Could not parse line: {line}")
            continue

    # 将卡片列表保存为JSON文件
    safe_name = sanitize_filename(pdf_filename)
    output_filename = f"{safe_name}_flashcards.json"
    with open(output_filename, "w", encoding='utf-8') as f:
        json.dump(flashcards_list, f, ensure_ascii=False, indent=4)
    
    print(f"Flashcards have been saved to {output_filename}")
    if MAX_FLASHCARDS > 0:
        print(f"Generated {len(flashcards_list)} flashcards (limit: {MAX_FLASHCARDS})")
    else:
        print(f"Generated {len(flashcards_list)} flashcards")
    # 直接在控制台输出JSON内容
    print("\n--- Generated JSON Output ---")
    print(json.dumps(flashcards_list, ensure_ascii=False, indent=4))


# Main script execution
if __name__ == "__main__":
    # 确保SOURCE_DOCUMENTS文件夹存在
    source_docs_path = os.path.join(ROOT_DIRECTORY, 'SOURCE_DOCUMENTS')
    os.makedirs(source_docs_path, exist_ok=True)
    
    # 获取SOURCE_DOCUMENTS文件夹中的第一个PDF、PPTX或PPT文件
    supported_files = [f for f in os.listdir(source_docs_path) 
                       if f.lower().endswith(('.pdf', '.pptx', '.ppt'))]
    
    if supported_files:
        filename = supported_files[0]
        file_path = os.path.join(source_docs_path, filename)
        print(f"Reading file: {filename}")
        
        if filename.lower().endswith('.pdf'):
            text_content = read_pdf(file_path)
        elif filename.lower().endswith('.pptx'):
            if not PPTX_AVAILABLE:
                print("Error: python-pptx not installed. Install with: pip install python-pptx")
                exit(1)
            text_content = read_pptx(file_path)
        elif filename.lower().endswith('.ppt'):
            text_content = read_ppt(file_path)
        
        create_anki_cards(text_content, filename)
    else:
        print("Error: No PDF, PPTX, or PPT files found in SOURCE_DOCUMENTS folder.")
        print("Please add a supported file to the SOURCE_DOCUMENTS folder.")
        text_content = "What is the capital of France?;Paris\nWhat is 2+2?;4"
        create_anki_cards(text_content)